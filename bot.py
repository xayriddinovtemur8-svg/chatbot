import os
import io
import json
import random
import string
import requests
import psycopg2
import fitz
import base64
from datetime import datetime, timedelta, date
from groq import Groq
from dotenv import load_dotenv
from gtts import gTTS
from langdetect import detect
from telegram import Update, BotCommand, InlineKeyboardButton, InlineKeyboardMarkup, LabeledPrice
from telegram.ext import (
    ApplicationBuilder, MessageHandler, CommandHandler,
    CallbackQueryHandler, PreCheckoutQueryHandler, filters, ContextTypes
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocPt, RGBColor as DocRGB

load_dotenv()

TELEGRAM_TOKEN    = os.getenv("TELEGRAM_TOKEN")
GROQ_API_KEY      = os.getenv("GROQ_API_KEY")
DATABASE_URL      = os.getenv("DATABASE_URL")
CARD_NUMBER       = os.getenv("CARD_NUMBER")
WEATHER_API_KEY   = os.getenv("WEATHER_API_KEY")
COINGECKO_API_KEY = os.getenv("COINGECKO_API_KEY")

ADMIN_ID       = 8230883785
ADMIN_USERNAME = "temur_uzb7779"
BOT_NAME       = "ChatBot Pro"

STARS_STANDARD = 770
STARS_PREMIUM  = 1150
STARS_GROUP    = 8547
STARS_IMAGINE  = 50
USDT_STANDARD  = 10
USDT_PREMIUM   = 15

COINS_PER_MSG  = 5
COINS_STANDARD = 10000
COINS_PREMIUM  = 100000
AFFILIATE_PCT  = 25
MAX_WARNINGS   = 10

BAD_WORDS = [
    "ублюдок","сука","пизда","хуй","мудак","залупа","пиздец","блядь","ёбаный",
    "fuck","shit","asshole","bitch","bastard","dick","cunt",
    "ahmoq","tentak","yaramas","eshak","haromzoda","onangni","otangni","egangni","cho'chqa"
]

SEARCH_KEYWORDS = [
    "today","now","current","latest","news","price","weather","rate",
    "bugun","hozir","narx","kurs","yangilik","ob-havo","oxirgi",
    "сегодня","сейчас","курс","цена","новости","погода"
]

GTTS_LANG_MAP = {
    "uz":"ru","en":"en","ru":"ru","tr":"tr","de":"de","fr":"fr",
    "es":"es","ar":"ar","ko":"ko","ja":"ja","zh-cn":"zh-CN","it":"it","pt":"pt","hi":"hi"
}

LANGUAGES = {
    "en":("🇬🇧","English"),  "ru":("🇷🇺","Русский"),  "uz":("🇺🇿","O'zbek"),
    "tr":("🇹🇷","Türkçe"),   "de":("🇩🇪","Deutsch"),  "fr":("🇫🇷","Français"),
    "es":("🇪🇸","Español"),  "ar":("🇸🇦","العربية"),  "ko":("🇰🇷","한국어"),
    "ja":("🇯🇵","日本語"),   "zh":("🇨🇳","中文"),      "it":("🇮🇹","Italiano"),
    "pt":("🇵🇹","Português"),"hi":("🇮🇳","हिंदी"),
}

TEXTS = {
"en":{
"welcome":(
    "👋 *Welcome to {bot}!*\n━━━━━━━━━━━━━━━━━━━━\n"
    "📋 *Plan:* {pe} {plan}\n🪙 *Coins:* {coins:,}\n\n"
    "🤖 *AI Features:*\n"
    "💬 Chat • 🌐 Search • 📄 PDF • 🖼️ Images • 🎤 Voice\n"
    "🎨 AI Image • 🌐 Translate • 💻 Code • 📋 Document\n\n"
    "📦 *Create Files:*\n"
    "📊 /pptx • 📝 /word • 👤 /cv • 📧 /email • 📱 /post\n\n"
    "🌤 *Real-time:* /weather • /crypto • /news\n\n"
    "🎁 /coins • /referral • /affiliate • /gift • /top • /stats\n"
    "⚙️ /updateplan • /language • /help"
),
"help":(
    "📌 *Commands*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "💬 Chat freely • 📄 Send PDF • 🖼️ Send photo • 🎤 Voice\n"
    "/translate /code /document ⭐\n"
    "/pptx 💎 /word 💎 /cv ⭐ /email ⭐ /post /biznes\n"
    "/imagine 💎 /ai_sound ⭐\n\n"
    "🌤 /weather /crypto /news\n\n"
    "💰 /updateplan /coins /referral /affiliate\n"
    "/gift /top /stats /promo /language /reset\n\n"
    "⭐ = Standard+  💎 = Premium only"
),
"plan_free":"🆓 FREE","plan_standard":"⭐ STANDARD",
"plan_premium":"💎 PREMIUM","plan_admin":"👑 ADMIN",
"limit_reached":"❌ *Daily limit reached!*\nUpgrade with /updateplan",
"blocked":"🚫 You are blocked. Contact @{admin}",
"cleared":"✅ Chat history cleared!",
"gen_voice":"⏳ Generating voice...","gen_image":"⏳ Generating image...",
"translating":"⏳ Translating...","writing_code":"⏳ Writing code...",
"writing_doc":"⏳ Creating document...","writing_cv":"⏳ Writing CV...",
"writing_email":"⏳ Writing email...","writing_post":"⏳ Writing post...",
"writing_biz":"⏳ Writing business plan...",
"creating_pptx":"⏳ Creating *{tp}*...","creating_word":"⏳ Creating *{tp}*...",
"ready":"✅ Ready!",
"locked_std":"⭐ *Requires Standard or Premium!*\n/updateplan",
"locked_prm":"💎 *Requires Premium!*\n/updateplan",
"locked_img":"💎 AI Image requires *Premium!*\nOr pay *{s}⭐* per image.",
"ex_voice":"Example: `/ai_sound Hello world`",
"ex_imagine":"Example: `/imagine beautiful sunset`",
"ex_translate":"Example: `/translate Hello → Uzbek`",
"ex_code":"Example: `/code Python fibonacci`",
"ex_document":"Example: `/document job application`",
"ex_cv":"Example: `/cv Python dev, 3 years`",
"ex_email":"Example: `/email follow up after interview`",
"ex_post":"Example: `/post grand opening coffee shop`",
"ex_biznes":"Example: `/biznes online clothing store`",
"ex_pptx":"Example: `/pptx artificial intelligence`",
"ex_word":"Example: `/word business proposal`",
"ex_weather":"Example: `/weather London`",
"ex_crypto":"Example: `/crypto bitcoin`",
"weather_err":"❌ City not found. Example: `/weather London`",
"crypto_err":"❌ Coin not found. Example: `/crypto bitcoin`",
"choose_lang":"🌐 *Choose your language:*",
"lang_cooldown":"⏳ You can change language once per 24 hours.",
"pdf_only":"📄 Please send a PDF file!",
"reading_pdf":"⏳ Reading PDF...",
"you_said":"🎤 *You said:* ",
"coins_info":(
    "🪙 *Your Coins: {coins:,}*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "📈 *Earn:* +{cpm} per message • +{aff}% referrals • +25 daily\n\n"
    "🛒 *Spend:*\n• {std:,} coins = ⭐ Standard 30 days\n"
    "• {prm:,} coins = 💎 Premium 30 days"
),
"not_enough":"❌ Need *{need:,}* more coins.",
"coins_ok_std":"✅ *Standard* activated for 30 days! 🎉",
"coins_ok_prm":"✅ *Premium* activated for 30 days! 🎉",
"referral_info":(
    "👥 *Referral Program*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "🔗 *Your link:*\n`{link}`\n\n"
    "📊 Invited: *{count}* friends\n\n"
    "🎁 10 friends → ⭐ Standard 15 days — {ss}\n"
    "🎁 30 friends → 💎 Premium 15 days — {ps}\n\n"
    "🪙 Earn *{aff}%* of every referral's coins!"
),
"claim_std":"🎁 Claim Standard 15 days",
"claim_prm":"🎁 Claim Premium 15 days",
"claimed_std":"🎉 ⭐ *Standard* activated for 15 days!",
"claimed_prm":"🎉 💎 *Premium* activated for 15 days!",
"claim_err":"❌ Already claimed or not enough referrals!",
"aff_info":(
    "🤝 *Affiliate Program*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "🔗 *Your link:*\n`{link}`\n\n"
    "📊 Referrals: *{count}* | 🪙 Earned: *{earned:,}* coins\n\n"
    "Share your link and earn *{aff}%* of their coins!"
),
"stats_info":(
    "📊 *Your Statistics*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "{pe} *{plan}* | 📅 {exp}\n"
    "🪙 {coins:,} | 👥 {refs} | 🔥 {streak}d | 💬 {total:,}\n\n"
    "📈 *Today:*\n"
    "💬{chat} 🌐{search} 🖼️{image} 📄{pdf}\n"
    "👤{cv} 📧{email} 🔊{tts} 📱{post}"
),
"gift_done":"🎁 *Already claimed today!*\nCome back tomorrow 😊",
"gift_ok":"🎁 *Daily Bonus!*\nStreak: *{streak} days* 🔥\n+5 messages & +25 coins!",
"gift_streak":"🎉 *10 day streak!*\n⭐ Standard activated for 10 days!",
"top_title":"🏆 *Top 10 Most Active Users*\n━━━━━━━━━━━━━━━━━━━━\n\n",
"promo_enter":"Enter your code: `/promo CODE`",
"promo_bad":"❌ Invalid or expired promo code!",
"promo_used":"❌ You already used this code!",
"promo_ok":"✅ Promo code applied! *{reward}*",
"group_info":(
    "👥 *Group Mode — {s}⭐ ($111.11)*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "✅ All group members can use the bot\n"
    "✅ Type `bot [question]` in group\n"
    "✅ Bot replies in the group chat\n\n"
    "After payment → add bot to group → /activate_group"
),
"group_ok":"✅ *Bot activated for this group!*\nMembers can type `bot [question]`",
"warning":"⚠️ *Warning {c}/{m}:* Please be respectful!",
"banned":"🚫 You have been banned for repeated abuse.",
"achievement":"🏆 *Achievement: 100,000 messages!*\n⭐ Standard activated for 5 days!",
"maintenance":"🔧 *Bot is temporarily down.*\nWe'll be back shortly! Sorry for the inconvenience.",
"updateplan_txt":(
    "💰 *Plans & Pricing*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "🆓 *FREE*\n• Chat 20/day • Translate • Code\n\n"
    "⭐ *STANDARD — {us}$ / {ss}⭐ / {cs:,}🪙*\n"
    "• 30/day + PDF • CV • Email • AI Voice • Document\n\n"
    "💎 *PREMIUM — {up}$ / {sp}⭐ / {cp:,}🪙*\n"
    "• Unlimited + AI Image • PowerPoint • Word\n\n"
    "👥 *GROUP — {sg}⭐ ($111.11)*\n"
    "• Bot for your entire group\n\n"
    "━━━━━━━━━━━━━━━━━━━━\n"
    "{pe} *{plan}* | 📅 {exp} | 🪙 {coins:,}"
),
"pay_card":(
    "💳 *Pay to card:*\n`{card}`\n\n"
    "💵 Amount: *{amt} USDT*\n\n"
    "1. Take a screenshot of payment\n"
    "2. Send to @{admin}\n"
    "3. Activated within 1 hour ✅"
),
"btn_ss":f"⭐ Standard — {STARS_STANDARD}⭐",
"btn_ps":f"💎 Premium — {STARS_PREMIUM}⭐",
"btn_sc":f"⭐ Standard — {COINS_STANDARD:,}🪙",
"btn_pc":f"💎 Premium — {COINS_PREMIUM:,}🪙",
"btn_su":f"💳 Standard — {USDT_STANDARD}$",
"btn_pu":f"💳 Premium — {USDT_PREMIUM}$",
"btn_grp":f"👥 Group — {STARS_GROUP}⭐",
"btn_ustd":"⭐ Upgrade to Standard",
"btn_uprm":"💎 Upgrade to Premium",
"btn_contact":"💬 Contact Admin",
"btn_img":f"🎨 Pay {STARS_IMAGINE}⭐ per image",
},
"uz":{
"welcome":(
    "👋 *{bot} ga xush kelibsiz!*\n━━━━━━━━━━━━━━━━━━━━\n"
    "📋 *Tarif:* {pe} {plan}\n🪙 *Coinlar:* {coins:,}\n\n"
    "🤖 *AI Imkoniyatlar:*\n"
    "💬 Suhbat • 🌐 Qidiruv • 📄 PDF • 🖼️ Rasm • 🎤 Ovoz\n"
    "🎨 AI Rasm • 🌐 Tarjima • 💻 Kod • 📋 Hujjat\n\n"
    "📦 *Fayl Yaratish:*\n"
    "📊 /pptx • 📝 /word • 👤 /cv • 📧 /email • 📱 /post\n\n"
    "🌤 *Real-vaqt:* /weather • /crypto • /news\n\n"
    "🎁 /coins • /referral • /affiliate • /gift • /top • /stats\n"
    "⚙️ /updateplan • /language • /help"
),
"help":(
    "📌 *Buyruqlar*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "💬 Suhbat • 📄 PDF • 🖼️ Rasm • 🎤 Ovoz\n"
    "/translate /code /document ⭐\n"
    "/pptx 💎 /word 💎 /cv ⭐ /email ⭐ /post /biznes\n"
    "/imagine 💎 /ai_sound ⭐\n\n"
    "🌤 /weather /crypto /news\n\n"
    "💰 /updateplan /coins /referral /affiliate\n"
    "/gift /top /stats /promo /language /reset\n\n"
    "⭐ = Standart+  💎 = Faqat Premium"
),
"plan_free":"🆓 BEPUL","plan_standard":"⭐ STANDART",
"plan_premium":"💎 PREMIUM","plan_admin":"👑 ADMIN",
"limit_reached":"❌ *Kunlik limit tugadi!*\n/updateplan orqali yangilang",
"blocked":"🚫 Bloklangansiz. @{admin} bilan bog'laning",
"cleared":"✅ Suhbat tarixi tozalandi!",
"gen_voice":"⏳ Ovoz yaratilmoqda...","gen_image":"⏳ Rasm yaratilmoqda...",
"translating":"⏳ Tarjima qilinmoqda...","writing_code":"⏳ Kod yozilmoqda...",
"writing_doc":"⏳ Hujjat yaratilmoqda...","writing_cv":"⏳ CV yozilmoqda...",
"writing_email":"⏳ Email yozilmoqda...","writing_post":"⏳ Post yozilmoqda...",
"writing_biz":"⏳ Biznes reja yozilmoqda...",
"creating_pptx":"⏳ *{tp}* yaratilmoqda...","creating_word":"⏳ *{tp}* yaratilmoqda...",
"ready":"✅ Tayyor!",
"locked_std":"⭐ *Standart yoki Premium kerak!*\n/updateplan",
"locked_prm":"💎 *Premium kerak!*\n/updateplan",
"locked_img":"💎 AI Rasm *Premium* talab qiladi!\nYoki *{s}⭐* to'lang.",
"ex_voice":"Misol: `/ai_sound Salom dunyo`",
"ex_imagine":"Misol: `/imagine tog' quyosh botishi`",
"ex_translate":"Misol: `/translate Hello → O'zbek`",
"ex_code":"Misol: `/code Python fibonacci`",
"ex_document":"Misol: `/document ish arizasi`",
"ex_cv":"Misol: `/cv Python dasturchi, 3 yil`",
"ex_email":"Misol: `/email intervyudan keyin`",
"ex_post":"Misol: `/post kofe do'kon ochilishi`",
"ex_biznes":"Misol: `/biznes online kiyim do'koni`",
"ex_pptx":"Misol: `/pptx sun'iy intellekt`",
"ex_word":"Misol: `/word biznes taklif`",
"ex_weather":"Misol: `/weather Toshkent`",
"ex_crypto":"Misol: `/crypto bitcoin`",
"weather_err":"❌ Shahar topilmadi. Misol: `/weather Toshkent`",
"crypto_err":"❌ Topilmadi. Misol: `/crypto bitcoin`",
"choose_lang":"🌐 *Tilni tanlang:*",
"lang_cooldown":"⏳ Tilni 24 soatda 1 marta o'zgartirish mumkin.",
"pdf_only":"📄 Iltimos PDF fayl yuboring!",
"reading_pdf":"⏳ PDF o'qilmoqda...",
"you_said":"🎤 *Siz dedingiz:* ",
"coins_info":(
    "🪙 *Coinlaringiz: {coins:,}*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "📈 *Yig'ish:* +{cpm} xabar • +{aff}% referal • +25 kunlik\n\n"
    "🛒 *Sarflash:*\n• {std:,} coin = ⭐ Standart 30 kun\n"
    "• {prm:,} coin = 💎 Premium 30 kun"
),
"not_enough":"❌ Yana *{need:,}* coin kerak.",
"coins_ok_std":"✅ *Standart* 30 kunga yoqildi! 🎉",
"coins_ok_prm":"✅ *Premium* 30 kunga yoqildi! 🎉",
"referral_info":(
    "👥 *Referal Dasturi*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "🔗 *Sizning havola:*\n`{link}`\n\n"
    "📊 Taklif qilingan: *{count}* do'st\n\n"
    "🎁 10 do'st → ⭐ Standart 15 kun — {ss}\n"
    "🎁 30 do'st → 💎 Premium 15 kun — {ps}\n\n"
    "🪙 Har bir referaldan *{aff}%* coin olasiz!"
),
"claim_std":"🎁 Standart 15 kun olish",
"claim_prm":"🎁 Premium 15 kun olish",
"claimed_std":"🎉 ⭐ *Standart* 15 kunga yoqildi!",
"claimed_prm":"🎉 💎 *Premium* 15 kunga yoqildi!",
"claim_err":"❌ Allaqachon olingan yoki referal yetarli emas!",
"aff_info":(
    "🤝 *Affiliate Dasturi*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "🔗 *Sizning havola:*\n`{link}`\n\n"
    "📊 Referallar: *{count}* | 🪙 Ishlangan: *{earned:,}* coin\n\n"
    "Havolangizni ulashing va *{aff}%* coin oling!"
),
"stats_info":(
    "📊 *Statistikangiz*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "{pe} *{plan}* | 📅 {exp}\n"
    "🪙 {coins:,} | 👥 {refs} | 🔥 {streak}k | 💬 {total:,}\n\n"
    "📈 *Bugun:*\n"
    "💬{chat} 🌐{search} 🖼️{image} 📄{pdf}\n"
    "👤{cv} 📧{email} 🔊{tts} 📱{post}"
),
"gift_done":"🎁 *Bugun bonus allaqachon olindi!*\nErtaga qaytib keling 😊",
"gift_ok":"🎁 *Kunlik Bonus!*\nStreak: *{streak} kun* 🔥\n+5 xabar & +25 coin!",
"gift_streak":"🎉 *10 kun streak!*\n⭐ Standart 10 kunga yoqildi!",
"top_title":"🏆 *Top 10 Eng Faol Foydalanuvchilar*\n━━━━━━━━━━━━━━━━━━━━\n\n",
"promo_enter":"Kodni kiriting: `/promo KOD`",
"promo_bad":"❌ Noto'g'ri yoki muddati o'tgan promo kod!",
"promo_used":"❌ Bu kodni allaqachon ishlatgansiz!",
"promo_ok":"✅ Promo kod qo'llandi! *{reward}*",
"group_info":(
    "👥 *Guruh Rejimi — {s}⭐ ($111.11)*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "✅ Barcha a'zolar botdan foydalana oladi\n"
    "✅ Guruhda `bot [savol]` deb yozing\n"
    "✅ Bot guruhda javob beradi\n\n"
    "To'lovdan keyin → botni guruhga qo'shing → /activate_group"
),
"group_ok":"✅ *Bot guruh uchun yoqildi!*\nA'zolar `bot [savol]` deb yoza oladi",
"warning":"⚠️ *Ogohlantirish {c}/{m}:* Iltimos hurmat bilan muomala qiling!",
"banned":"🚫 Qayta-qayta haqorat qilgani uchun bloklandi.",
"achievement":"🏆 *Yutuq: 100,000 xabar!*\n⭐ Standart 5 kunga yoqildi!",
"maintenance":"🔧 *Bot vaqtinchalik ishlamaydi.*\nTez orada qayta ishlaydi! Uzr so'raymiz.",
"updateplan_txt":(
    "💰 *Tariflar va Narxlar*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "🆓 *BEPUL*\n• Suhbat 20/kun • Tarjima • Kod\n\n"
    "⭐ *STANDART — {us}$ / {ss}⭐ / {cs:,}🪙*\n"
    "• 30/kun + PDF • CV • Email • AI Ovoz • Hujjat\n\n"
    "💎 *PREMIUM — {up}$ / {sp}⭐ / {cp:,}🪙*\n"
    "• Cheksiz + AI Rasm • PowerPoint • Word\n\n"
    "👥 *GURUH — {sg}⭐ ($111.11)*\n"
    "• Butun guruh uchun bot\n\n"
    "━━━━━━━━━━━━━━━━━━━━\n"
    "{pe} *{plan}* | 📅 {exp} | 🪙 {coins:,}"
),
"pay_card":(
    "💳 *Kartaga to'lang:*\n`{card}`\n\n"
    "💵 Miqdor: *{amt} USDT*\n\n"
    "1. To'lov skrinshoti\n"
    "2. @{admin} ga yuboring\n"
    "3. 1 soat ichida yoqiladi ✅"
),
"btn_ss":f"⭐ Standart — {STARS_STANDARD}⭐",
"btn_ps":f"💎 Premium — {STARS_PREMIUM}⭐",
"btn_sc":f"⭐ Standart — {COINS_STANDARD:,}🪙",
"btn_pc":f"💎 Premium — {COINS_PREMIUM:,}🪙",
"btn_su":f"💳 Standart — {USDT_STANDARD}$",
"btn_pu":f"💳 Premium — {USDT_PREMIUM}$",
"btn_grp":f"👥 Guruh — {STARS_GROUP}⭐",
"btn_ustd":"⭐ Standartga o'tish",
"btn_uprm":"💎 Premiumga o'tish",
"btn_contact":"💬 Admin bilan bog'lanish",
"btn_img":f"🎨 {STARS_IMAGINE}⭐ to'lab rasm",
},
"ru":{
"welcome":(
    "👋 *Добро пожаловать в {bot}!*\n━━━━━━━━━━━━━━━━━━━━\n"
    "📋 *Тариф:* {pe} {plan}\n🪙 *Монеты:* {coins:,}\n\n"
    "🤖 *AI Возможности:*\n"
    "💬 Чат • 🌐 Поиск • 📄 PDF • 🖼️ Фото • 🎤 Голос\n"
    "🎨 AI Фото • 🌐 Перевод • 💻 Код • 📋 Документ\n\n"
    "📦 *Создание Файлов:*\n"
    "📊 /pptx • 📝 /word • 👤 /cv • 📧 /email • 📱 /post\n\n"
    "🌤 *Реальное время:* /weather • /crypto • /news\n\n"
    "🎁 /coins • /referral • /affiliate • /gift • /top • /stats\n"
    "⚙️ /updateplan • /language • /help"
),
"help":(
    "📌 *Команды*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "💬 Чат • 📄 PDF • 🖼️ Фото • 🎤 Голос\n"
    "/translate /code /document ⭐\n"
    "/pptx 💎 /word 💎 /cv ⭐ /email ⭐ /post /biznes\n"
    "/imagine 💎 /ai_sound ⭐\n\n"
    "🌤 /weather /crypto /news\n\n"
    "💰 /updateplan /coins /referral /affiliate\n"
    "/gift /top /stats /promo /language /reset\n\n"
    "⭐ = Standard+  💎 = Только Premium"
),
"plan_free":"🆓 БЕСПЛАТНО","plan_standard":"⭐ STANDARD",
"plan_premium":"💎 PREMIUM","plan_admin":"👑 ADMIN",
"limit_reached":"❌ *Дневной лимит исчерпан!*\n/updateplan",
"blocked":"🚫 Вы заблокированы. @{admin}",
"cleared":"✅ История очищена!",
"gen_voice":"⏳ Генерация голоса...","gen_image":"⏳ Генерация изображения...",
"translating":"⏳ Перевожу...","writing_code":"⏳ Пишу код...",
"writing_doc":"⏳ Создаю документ...","writing_cv":"⏳ Пишу резюме...",
"writing_email":"⏳ Пишу email...","writing_post":"⏳ Пишу пост...",
"writing_biz":"⏳ Пишу бизнес-план...",
"creating_pptx":"⏳ Создаю *{tp}*...","creating_word":"⏳ Создаю *{tp}*...",
"ready":"✅ Готово!",
"locked_std":"⭐ *Требуется Standard или Premium!*\n/updateplan",
"locked_prm":"💎 *Требуется Premium!*\n/updateplan",
"locked_img":"💎 AI Фото требует *Premium!*\nИли *{s}⭐* за изображение.",
"ex_voice":"Пример: `/ai_sound Привет мир`",
"ex_imagine":"Пример: `/imagine красивый закат`",
"ex_translate":"Пример: `/translate Hello → Русский`",
"ex_code":"Пример: `/code Python fibonacci`",
"ex_document":"Пример: `/document заявление на работу`",
"ex_cv":"Пример: `/cv Python разработчик, 3 года`",
"ex_email":"Пример: `/email после собеседования`",
"ex_post":"Пример: `/post открытие кофейни`",
"ex_biznes":"Пример: `/biznes онлайн магазин одежды`",
"ex_pptx":"Пример: `/pptx искусственный интеллект`",
"ex_word":"Пример: `/word бизнес-предложение`",
"ex_weather":"Пример: `/weather Москва`",
"ex_crypto":"Пример: `/crypto bitcoin`",
"weather_err":"❌ Город не найден. Пример: `/weather Москва`",
"crypto_err":"❌ Не найдено. Пример: `/crypto bitcoin`",
"choose_lang":"🌐 *Выберите язык:*",
"lang_cooldown":"⏳ Язык можно менять раз в 24 часа.",
"pdf_only":"📄 Отправьте PDF файл!",
"reading_pdf":"⏳ Читаю PDF...",
"you_said":"🎤 *Вы сказали:* ",
"coins_info":(
    "🪙 *Ваши Монеты: {coins:,}*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "📈 *Заработать:* +{cpm} за сообщение • +{aff}% рефералы • +25 ежедневно\n\n"
    "🛒 *Потратить:*\n• {std:,} монет = ⭐ Standard 30 дней\n"
    "• {prm:,} монет = 💎 Premium 30 дней"
),
"not_enough":"❌ Нужно ещё *{need:,}* монет.",
"coins_ok_std":"✅ *Standard* активирован на 30 дней! 🎉",
"coins_ok_prm":"✅ *Premium* активирован на 30 дней! 🎉",
"referral_info":(
    "👥 *Реферальная Программа*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "🔗 *Ваша ссылка:*\n`{link}`\n\n"
    "📊 Приглашено: *{count}* друзей\n\n"
    "🎁 10 друзей → ⭐ Standard 15 дней — {ss}\n"
    "🎁 30 друзей → 💎 Premium 15 дней — {ps}\n\n"
    "🪙 Зарабатывайте *{aff}%* монет каждого реферала!"
),
"claim_std":"🎁 Получить Standard 15 дней",
"claim_prm":"🎁 Получить Premium 15 дней",
"claimed_std":"🎉 ⭐ *Standard* активирован на 15 дней!",
"claimed_prm":"🎉 💎 *Premium* активирован на 15 дней!",
"claim_err":"❌ Уже получено или недостаточно рефералов!",
"aff_info":(
    "🤝 *Аффилиат Программа*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "🔗 *Ваша ссылка:*\n`{link}`\n\n"
    "📊 Рефералов: *{count}* | 🪙 Заработано: *{earned:,}* монет\n\n"
    "Делитесь ссылкой и получайте *{aff}%* их монет!"
),
"stats_info":(
    "📊 *Ваша Статистика*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "{pe} *{plan}* | 📅 {exp}\n"
    "🪙 {coins:,} | 👥 {refs} | 🔥 {streak}д | 💬 {total:,}\n\n"
    "📈 *Сегодня:*\n"
    "💬{chat} 🌐{search} 🖼️{image} 📄{pdf}\n"
    "👤{cv} 📧{email} 🔊{tts} 📱{post}"
),
"gift_done":"🎁 *Бонус уже получен сегодня!*\nВозвращайтесь завтра 😊",
"gift_ok":"🎁 *Ежедневный Бонус!*\nСерия: *{streak} дней* 🔥\n+5 сообщений & +25 монет!",
"gift_streak":"🎉 *10 дней подряд!*\n⭐ Standard активирован на 10 дней!",
"top_title":"🏆 *Топ 10 Самых Активных*\n━━━━━━━━━━━━━━━━━━━━\n\n",
"promo_enter":"Введите код: `/promo КОД`",
"promo_bad":"❌ Неверный или истёкший промокод!",
"promo_used":"❌ Вы уже использовали этот код!",
"promo_ok":"✅ Промокод применён! *{reward}*",
"group_info":(
    "👥 *Групповой Режим — {s}⭐ ($111.11)*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "✅ Все участники группы используют бота\n"
    "✅ Пишите `bot [вопрос]` в группе\n"
    "✅ Бот отвечает прямо в чате\n\n"
    "После оплаты → добавьте бота → /activate_group"
),
"group_ok":"✅ *Бот активирован для группы!*\nУчастники могут писать `bot [вопрос]`",
"warning":"⚠️ *Предупреждение {c}/{m}:* Пожалуйста, будьте вежливы!",
"banned":"🚫 Заблокированы за повторные оскорбления.",
"achievement":"🏆 *Достижение: 100,000 сообщений!*\n⭐ Standard активирован на 5 дней!",
"maintenance":"🔧 *Бот временно недоступен.*\nСкоро вернётся! Приносим извинения.",
"updateplan_txt":(
    "💰 *Тарифы и Цены*\n━━━━━━━━━━━━━━━━━━━━\n\n"
    "🆓 *БЕСПЛАТНО*\n• Чат 20/день • Перевод • Код\n\n"
    "⭐ *STANDARD — {us}$ / {ss}⭐ / {cs:,}🪙*\n"
    "• 30/день + PDF • CV • Email • AI Голос • Документ\n\n"
    "💎 *PREMIUM — {up}$ / {sp}⭐ / {cp:,}🪙*\n"
    "• Безлимит + AI Фото • PowerPoint • Word\n\n"
    "👥 *ГРУППА — {sg}⭐ ($111.11)*\n"
    "• Бот для всей группы\n\n"
    "━━━━━━━━━━━━━━━━━━━━\n"
    "{pe} *{plan}* | 📅 {exp} | 🪙 {coins:,}"
),
"pay_card":(
    "💳 *Оплатите на карту:*\n`{card}`\n\n"
    "💵 Сумма: *{amt} USDT*\n\n"
    "1. Скриншот оплаты\n"
    "2. Отправьте @{admin}\n"
    "3. Активация в течение 1 часа ✅"
),
"btn_ss":f"⭐ Standard — {STARS_STANDARD}⭐",
"btn_ps":f"💎 Premium — {STARS_PREMIUM}⭐",
"btn_sc":f"⭐ Standard — {COINS_STANDARD:,}🪙",
"btn_pc":f"💎 Premium — {COINS_PREMIUM:,}🪙",
"btn_su":f"💳 Standard — {USDT_STANDARD}$",
"btn_pu":f"💳 Premium — {USDT_PREMIUM}$",
"btn_grp":f"👥 Группа — {STARS_GROUP}⭐",
"btn_ustd":"⭐ Перейти на Standard",
"btn_uprm":"💎 Перейти на Premium",
"btn_contact":"💬 Написать Админу",
"btn_img":f"🎨 {STARS_IMAGINE}⭐ за изображение",
},
}

for _l in LANGUAGES:
    if _l not in TEXTS: TEXTS[_l] = TEXTS["en"]

def tx(lang, key, **kw):
    s = TEXTS.get(lang, TEXTS["en"]).get(key, TEXTS["en"].get(key, ""))
    if kw:
        try: s = s.format(**kw)
        except: pass
    return s

# Groq AI client
ai = Groq(api_key=GROQ_API_KEY)
user_histories = {}

def db():
    return psycopg2.connect(DATABASE_URL)

def init_db():
    conn = db(); c = conn.cursor()
    c.execute("""CREATE TABLE IF NOT EXISTS users (
        id SERIAL PRIMARY KEY, user_id BIGINT UNIQUE NOT NULL,
        username TEXT, full_name TEXT, plan TEXT DEFAULT 'free',
        expires_at TIMESTAMP, is_blocked BOOLEAN DEFAULT FALSE,
        is_admin BOOLEAN DEFAULT FALSE, joined_at TIMESTAMP DEFAULT NOW(),
        language TEXT DEFAULT 'en', language_changed_at TIMESTAMP,
        coins INTEGER DEFAULT 0, warning_count INTEGER DEFAULT 0,
        referral_code TEXT, affiliate_code TEXT, referred_by BIGINT,
        referral_count INTEGER DEFAULT 0, affiliate_earnings INTEGER DEFAULT 0,
        claimed_standard BOOLEAN DEFAULT FALSE, claimed_premium BOOLEAN DEFAULT FALSE,
        streak_count INTEGER DEFAULT 0, last_gift_claim DATE,
        streak_claimed BOOLEAN DEFAULT FALSE, group_plan TEXT DEFAULT 'none',
        total_messages INTEGER DEFAULT 0, achievement_100k BOOLEAN DEFAULT FALSE,
        usage_chat INTEGER DEFAULT 0, usage_search INTEGER DEFAULT 0,
        usage_image INTEGER DEFAULT 0, usage_post INTEGER DEFAULT 0,
        usage_biznes INTEGER DEFAULT 0, usage_pdf INTEGER DEFAULT 0,
        usage_cv INTEGER DEFAULT 0, usage_email INTEGER DEFAULT 0,
        usage_tts INTEGER DEFAULT 0, last_reset DATE DEFAULT CURRENT_DATE
    )""")
    c.execute("""CREATE TABLE IF NOT EXISTS user_memory (
        user_id BIGINT PRIMARY KEY, name TEXT, facts TEXT
    )""")
    c.execute("""CREATE TABLE IF NOT EXISTS group_modes (
        chat_id BIGINT PRIMARY KEY, activated_by BIGINT,
        activated_at TIMESTAMP DEFAULT NOW()
    )""")
    c.execute("""CREATE TABLE IF NOT EXISTS promo_codes (
        id SERIAL PRIMARY KEY, code TEXT UNIQUE,
        reward_type TEXT, reward_value INTEGER,
        max_uses INTEGER DEFAULT 1, used_count INTEGER DEFAULT 0,
        expires_at TIMESTAMP, created_at TIMESTAMP DEFAULT NOW()
    )""")
    c.execute("""CREATE TABLE IF NOT EXISTS promo_uses (
        user_id BIGINT, code TEXT, PRIMARY KEY(user_id, code)
    )""")
    c.execute("UPDATE users SET is_admin=TRUE WHERE user_id=%s", (ADMIN_ID,))
    conn.commit(); conn.close()

def rnd(n=8):
    return ''.join(random.choices(string.ascii_uppercase + string.digits, k=n))

def ensure_user(uid, username=None, full_name=None, referred_by=None):
    try:
        conn = db(); c = conn.cursor()
        c.execute("""INSERT INTO users(user_id,username,full_name,referral_code,affiliate_code,referred_by,is_admin)
            VALUES(%s,%s,%s,%s,%s,%s,%s) ON CONFLICT(user_id) DO UPDATE SET
            username=EXCLUDED.username, full_name=EXCLUDED.full_name""",
            (uid, username, full_name, rnd(8), rnd(10), referred_by, uid == ADMIN_ID))
        conn.commit()
        if referred_by:
            c.execute("UPDATE users SET referral_count=referral_count+1 WHERE user_id=%s", (referred_by,))
            conn.commit()
        conn.close()
    except: pass

def get_user(uid):
    D = {"plan":"free","expires_at":None,"is_blocked":False,"is_admin":False,
         "full_name":None,"username":None,"language":"en","language_changed_at":None,
         "coins":0,"warning_count":0,"referral_code":None,"affiliate_code":None,
         "referred_by":None,"referral_count":0,"affiliate_earnings":0,
         "claimed_standard":False,"claimed_premium":False,"streak_count":0,
         "last_gift_claim":None,"streak_claimed":False,"group_plan":"none",
         "total_messages":0,"achievement_100k":False}
    try:
        conn = db(); c = conn.cursor()
        c.execute("""SELECT plan,expires_at,is_blocked,is_admin,full_name,username,
            language,language_changed_at,coins,warning_count,referral_code,affiliate_code,
            referred_by,referral_count,affiliate_earnings,claimed_standard,claimed_premium,
            streak_count,last_gift_claim,streak_claimed,group_plan,total_messages,achievement_100k
            FROM users WHERE user_id=%s""", (uid,))
        row = c.fetchone(); conn.close()
        if not row: return D
        keys = ["plan","expires_at","is_blocked","is_admin","full_name","username",
                "language","language_changed_at","coins","warning_count","referral_code",
                "affiliate_code","referred_by","referral_count","affiliate_earnings",
                "claimed_standard","claimed_premium","streak_count","last_gift_claim",
                "streak_claimed","group_plan","total_messages","achievement_100k"]
        u = dict(zip(keys, row))
        for k in ["is_blocked","is_admin","claimed_standard","claimed_premium","streak_claimed","achievement_100k"]:
            u[k] = bool(u.get(k))
        for k in ["coins","warning_count","referral_count","affiliate_earnings","streak_count","total_messages"]:
            u[k] = u.get(k) or 0
        u["language"] = u.get("language") or "en"
        u["group_plan"] = u.get("group_plan") or "none"
        if not u["is_admin"] and u["expires_at"] and datetime.now() > u["expires_at"] and u["plan"] != "free":
            set_plan(uid, "free", None); u["plan"] = "free"; u["expires_at"] = None
        return u
    except: return D

def set_plan(uid, plan, days=30):
    try:
        conn = db(); c = conn.cursor()
        exp = datetime.now() + timedelta(days=days) if days else None
        c.execute("UPDATE users SET plan=%s,expires_at=%s WHERE user_id=%s", (plan, exp, uid))
        conn.commit(); conn.close()
    except: pass

def set_blocked(uid, v):
    try:
        conn = db(); c = conn.cursor()
        c.execute("UPDATE users SET is_blocked=%s WHERE user_id=%s", (v, uid))
        conn.commit(); conn.close()
    except: pass

def set_language(uid, lang):
    try:
        conn = db(); c = conn.cursor()
        c.execute("UPDATE users SET language=%s,language_changed_at=%s WHERE user_id=%s",
                  (lang, datetime.now(), uid))
        conn.commit(); conn.close()
    except: pass

def add_coins(uid, n):
    try:
        conn = db(); c = conn.cursor()
        c.execute("UPDATE users SET coins=coins+%s WHERE user_id=%s", (n, uid))
        conn.commit(); conn.close()
    except: pass

def add_warning(uid):
    try:
        conn = db(); c = conn.cursor()
        c.execute("UPDATE users SET warning_count=warning_count+1 WHERE user_id=%s", (uid,))
        c.execute("SELECT warning_count FROM users WHERE user_id=%s", (uid,))
        n = c.fetchone()[0]; conn.commit(); conn.close(); return n
    except: return 0

def inc_msg(uid, is_admin=False):
    if is_admin: return 0
    try:
        conn = db(); c = conn.cursor()
        c.execute("UPDATE users SET total_messages=total_messages+1,coins=coins+%s WHERE user_id=%s",
                  (COINS_PER_MSG, uid))
        conn.commit()
        c.execute("SELECT total_messages,referred_by FROM users WHERE user_id=%s", (uid,))
        row = c.fetchone(); conn.close()
        if row:
            total, ref = row
            if ref:
                bonus = int(COINS_PER_MSG * AFFILIATE_PCT / 100)
                add_coins(ref, bonus)
                try:
                    conn2 = db(); c2 = conn2.cursor()
                    c2.execute("UPDATE users SET affiliate_earnings=affiliate_earnings+%s WHERE user_id=%s",
                               (bonus, ref))
                    conn2.commit(); conn2.close()
                except: pass
            return total
        return 0
    except: return 0

def check_limit(uid, feature, limit):
    if limit == -1: return True
    if limit == 0: return False
    try:
        conn = db(); c = conn.cursor()
        today = datetime.now().date()
        c.execute("SELECT last_reset FROM users WHERE user_id=%s", (uid,))
        row = c.fetchone()
        if row and row[0] < today:
            c.execute("""UPDATE users SET usage_chat=0,usage_search=0,usage_image=0,
                usage_post=0,usage_biznes=0,usage_pdf=0,usage_cv=0,usage_email=0,
                usage_tts=0,last_reset=%s WHERE user_id=%s""", (today, uid))
            conn.commit()
        c.execute(f"SELECT usage_{feature} FROM users WHERE user_id=%s", (uid,))
        row = c.fetchone(); usage = row[0] if row else 0
        if usage >= limit: conn.close(); return False
        c.execute(f"UPDATE users SET usage_{feature}=usage_{feature}+1 WHERE user_id=%s", (uid,))
        conn.commit(); conn.close(); return True
    except: return True

def get_limits(plan, is_admin=False):
    ALL = {k: -1 for k in ["chat","search","image","post","biznes","pdf","cv","email",
                            "voice","pptx","word","tts","imagine","translate","code","document"]}
    if is_admin or plan == "premium": return ALL
    if plan == "standard":
        return {"chat":30,"search":30,"image":30,"post":30,"biznes":30,"pdf":30,
                "cv":30,"email":30,"voice":0,"pptx":0,"word":0,"tts":30,
                "imagine":0,"translate":30,"code":30,"document":30}
    return {"chat":20,"search":20,"image":20,"post":20,"biznes":20,"pdf":0,
            "cv":0,"email":0,"voice":0,"pptx":0,"word":0,"tts":0,
            "imagine":0,"translate":20,"code":20,"document":0}

def get_memory(uid):
    try:
        conn = db(); c = conn.cursor()
        c.execute("SELECT name,facts FROM user_memory WHERE user_id=%s", (uid,))
        row = c.fetchone(); conn.close()
        return {"name": row[0], "facts": row[1]} if row else None
    except: return None

def save_memory(uid, name, facts):
    try:
        conn = db(); c = conn.cursor()
        c.execute("""INSERT INTO user_memory(user_id,name,facts) VALUES(%s,%s,%s)
            ON CONFLICT(user_id) DO UPDATE SET name=%s,facts=%s""",
            (uid, name, facts, name, facts))
        conn.commit(); conn.close()
    except: pass

def is_group_active(chat_id):
    try:
        conn = db(); c = conn.cursor()
        c.execute("SELECT 1 FROM group_modes WHERE chat_id=%s", (chat_id,))
        r = c.fetchone(); conn.close(); return r is not None
    except: return False

def activate_group(chat_id, uid):
    try:
        conn = db(); c = conn.cursor()
        c.execute("INSERT INTO group_modes(chat_id,activated_by) VALUES(%s,%s) ON CONFLICT DO NOTHING",
                  (chat_id, uid))
        conn.commit(); conn.close()
    except: pass

def plan_info(u):
    lang = u["language"]
    if u["is_admin"]: return "👑", tx(lang, "plan_admin")
    p = u["plan"]
    em = {"free":"🆓","standard":"⭐","premium":"💎"}.get(p, "🆓")
    kk = {"free":"plan_free","standard":"plan_standard","premium":"plan_premium"}.get(p, "plan_free")
    return em, tx(lang, kk)

def exp_str(u):
    if u["is_admin"]: return "∞"
    if u.get("expires_at"): return u["expires_at"].strftime("%d.%m.%Y")
    return "—"

def get_weather(city):
    try:
        r = requests.get(
            f"https://api.openweathermap.org/data/2.5/weather"
            f"?q={city}&appid={WEATHER_API_KEY}&units=metric",
            timeout=10).json()
        if r.get("cod") != 200: return None
        w = r["weather"][0]; m = r["main"]; wind = r["wind"]
        em = {"Clear":"☀️","Clouds":"☁️","Rain":"🌧️","Snow":"❄️",
              "Thunderstorm":"⛈️","Drizzle":"🌦️","Mist":"🌫️","Fog":"🌫️"}.get(w["main"],"🌤")
        return (f"🌍 *{r['name']}, {r['sys']['country']}*\n━━━━━━━━━━━━━━━━━━━━\n"
                f"{em} *{w['description'].capitalize()}*\n\n"
                f"🌡 *Temp:* {m['temp']:.1f}°C (feels {m['feels_like']:.1f}°C)\n"
                f"💧 *Humidity:* {m['humidity']}%\n"
                f"💨 *Wind:* {wind['speed']} m/s\n"
                f"🔼 *Max:* {m['temp_max']:.1f}°C | 🔽 *Min:* {m['temp_min']:.1f}°C")
    except: return None

def get_crypto(coin):
    try:
        h = {"x-cg-demo-api-key": COINGECKO_API_KEY} if COINGECKO_API_KEY else {}
        r = requests.get(
            f"https://api.coingecko.com/api/v3/coins/{coin.lower()}"
            f"?localization=false&tickers=false&community_data=false&developer_data=false",
            headers=h, timeout=10).json()
        if "error" in r: return None
        md = r["market_data"]
        price = md["current_price"]["usd"]
        c24 = md["price_change_percentage_24h"] or 0
        c7 = md["price_change_percentage_7d"] or 0
        mc = md["market_cap"]["usd"]
        h24 = md["high_24h"]["usd"]; l24 = md["low_24h"]["usd"]
        arrow = "📈" if c24 >= 0 else "📉"
        return (f"💰 *{r['name']} ({r['symbol'].upper()})*\n━━━━━━━━━━━━━━━━━━━━\n"
                f"💵 *Price:* ${price:,.4f}\n{arrow} *24h:* {c24:+.2f}%\n"
                f"📊 *7d:* {c7:+.2f}%\n"
                f"📈 *High:* ${h24:,.4f} | 📉 *Low:* ${l24:,.4f}\n"
                f"🏦 *Market Cap:* ${mc:,.0f}\n"
                f"🏆 *Rank:* #{r.get('market_cap_rank','N/A')}")
    except: return None

def gen_image(prompt):
    try:
        url = (f"https://image.pollinations.ai/prompt/"
               f"{requests.utils.quote(prompt)}"
               f"?width=1024&height=1024&nologo=true&enhance=true")
        r = requests.get(url, timeout=60)
        return r.content if r.status_code == 200 else None
    except: return None

def detect_gtts(text):
    try: return GTTS_LANG_MAP.get(detect(text), "en")
    except: return "en"

def is_bad(text):
    return any(w in text.lower() for w in BAD_WORDS)

def make_pptx(title, slides):
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = RGBColor(0x0F, 0x0F, 0x1A)
    tb = sl.shapes.title; tb.text = title
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tb.text_frame.paragraphs[0].font.size = Pt(40)
    tb.text_frame.paragraphs[0].font.bold = True
    for s in slides:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = RGBColor(0x0F, 0x0F, 0x1A)
        ti = sl.shapes.title; ti.text = s["title"]
        ti.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x89, 0xB4, 0xFA)
        ti.text_frame.paragraphs[0].font.size = Pt(28)
        ti.text_frame.paragraphs[0].font.bold = True
        tf = sl.placeholders[1].text_frame; tf.clear()
        for i, pt in enumerate(s["points"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸ {pt}"
            p.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)
            p.font.size = Pt(18)
    prs.save("pres.pptx"); return "pres.pptx"

def make_docx(title, sections):
    doc = Document()
    h = doc.add_heading(title, 0)
    if h.runs: h.runs[0].font.color.rgb = DocRGB(0x1E, 0x3A, 0x8A)
    for sec in sections:
        doc.add_heading(sec["title"], level=1)
        for pt in sec["points"]: doc.add_paragraph(pt, style="List Bullet")
        doc.add_paragraph()
    doc.save("doc.docx"); return "doc.docx"

async def ai_chat(uid, text, do_search=False, limits=None):
    mem = get_memory(uid)
    mc = f"User info: name={mem['name']}, facts={mem['facts']}. " if mem and (mem.get("name") or mem.get("facts")) else ""
    sys_p = (
        "You are ChatBot Pro — a professional AI assistant. "
        "CRITICAL RULE: Always reply in the EXACT same language as the user's message. "
        "If user writes in Uzbek → reply in Uzbek. "
        "If user writes in English → reply in English. "
        "If user writes in Russian → reply in Russian. "
        "Never mix languages. Never add unnecessary text. "
        "If asked who created you or who is your owner → reply: '@temur_uzb7779 created me.' "
        "Be helpful, accurate and professional. " + mc
    )
    if uid not in user_histories:
        user_histories[uid] = [{"role":"system","content":sys_p}]
    else:
        user_histories[uid][0]["content"] = sys_p
    content = text
    if do_search and limits and limits.get("search") != 0:
        if any(k in text.lower() for k in SEARCH_KEYWORDS):
            content = f"{text}\n[Note: Provide the most accurate and up-to-date answer you know]"
    user_histories[uid].append({"role":"user","content":content})
    r = ai.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=user_histories[uid],
        max_tokens=2048
    )
    reply = r.choices[0].message.content
    user_histories[uid].append({"role":"assistant","content":reply})
    return reply

async def ai_once(prompt):
    r = ai.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role":"user","content":prompt}],
        max_tokens=2048
    )
    return r.choices[0].message.content

async def update_mem(uid, user_text, reply):
    mem = get_memory(uid)
    cn = mem["name"] if mem else ""; cf = mem["facts"] if mem else ""
    try:
        r = await ai_once(
            f"Extract user's name and key facts from this conversation.\n"
            f"Current data: name='{cn}', facts='{cf}'\n"
            f"User said: {user_text}\n"
            f"Reply ONLY valid JSON, nothing else: "
            f'{{\"name\":\"name or empty\",\"facts\":\"short facts\"}}'
        )
        r = r.strip()
        if "```" in r:
            r = r.split("```")[1]
            if r.startswith("json"): r = r[4:]
        d = json.loads(r)
        save_memory(uid, d.get("name", cn), d.get("facts", cf))
    except: pass

async def broadcast_maintenance(app):
    try:
        conn = db(); c = conn.cursor()
        c.execute("SELECT user_id, language FROM users WHERE is_blocked=FALSE")
        rows = c.fetchall(); conn.close()
        for uid, lang in rows:
            try:
                await app.bot.send_message(
                    chat_id=uid,
                    text=tx(lang or "en", "maintenance"),
                    parse_mode="Markdown")
            except: pass
    except: pass

async def start(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ref = None
    if ctx.args:
        try:
            code = ctx.args[0]; conn = db(); c = conn.cursor()
            c.execute("SELECT user_id FROM users WHERE referral_code=%s OR affiliate_code=%s", (code, code))
            row = c.fetchone(); conn.close()
            if row and row[0] != user.id: ref = row[0]
        except: pass
    ensure_user(user.id, user.username, user.full_name, ref)
    u = get_user(user.id); lang = u["language"]
    pe, pn = plan_info(u)
    await update.message.reply_text(
        tx(lang, "welcome", bot=BOT_NAME, pe=pe, plan=pn, coins=u["coins"]),
        parse_mode="Markdown")

async def help_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id)
    await update.message.reply_text(tx(u["language"], "help"), parse_mode="Markdown")

async def reset_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id)
    user_histories[user.id] = []
    await update.message.reply_text(tx(u["language"], "cleared"), parse_mode="Markdown")

async def weather_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    city = " ".join(ctx.args) if ctx.args else ""
    if not city:
        await update.message.reply_text(tx(lang,"ex_weather"),parse_mode="Markdown"); return
    result = get_weather(city)
    await update.message.reply_text(
        result if result else tx(lang,"weather_err"), parse_mode="Markdown")

async def crypto_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    coin = ctx.args[0] if ctx.args else ""
    if not coin:
        await update.message.reply_text(tx(lang,"ex_crypto"),parse_mode="Markdown"); return
    result = get_crypto(coin)
    await update.message.reply_text(
        result if result else tx(lang,"crypto_err"), parse_mode="Markdown")

async def news_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    await update.message.reply_text("⏳ Fetching latest news...")
    result = await ai_once(
        "Give 5 important world news headlines with brief descriptions. "
        "Format each as: 📌 **Title**\nBrief description\n\n"
        "Be informative and professional.")
    if result:
        await update.message.reply_text(
            f"📰 *Latest News*\n━━━━━━━━━━━━━━━━━━━━\n\n{result}",
            parse_mode="Markdown")
    else:
        await update.message.reply_text("❌ Could not fetch news. Try again.")

async def language_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    lca = u["language_changed_at"]
    if lca and datetime.now() - lca < timedelta(hours=24):
        await update.message.reply_text(tx(lang,"lang_cooldown"),parse_mode="Markdown"); return
    kb = []; row = []
    for code, (flag, name) in LANGUAGES.items():
        row.append(InlineKeyboardButton(f"{flag} {name}", callback_data=f"lang_{code}"))
        if len(row) == 2: kb.append(row); row = []
    if row: kb.append(row)
    await update.message.reply_text(
        tx(lang,"choose_lang"), reply_markup=InlineKeyboardMarkup(kb), parse_mode="Markdown")

async def updateplan_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    pe, pn = plan_info(u); exp = exp_str(u)
    # Admin uchun ham to'liq ko'rinadi lekin to'lov tugmalari bor
    kb = [
        [InlineKeyboardButton(tx(lang,"btn_ss"),callback_data="stars_std"),
         InlineKeyboardButton(tx(lang,"btn_ps"),callback_data="stars_prm")],
        [InlineKeyboardButton(tx(lang,"btn_sc"),callback_data="coins_std"),
         InlineKeyboardButton(tx(lang,"btn_pc"),callback_data="coins_prm")],
        [InlineKeyboardButton(tx(lang,"btn_su"),callback_data="card_std"),
         InlineKeyboardButton(tx(lang,"btn_pu"),callback_data="card_prm")],
        [InlineKeyboardButton(tx(lang,"btn_grp"),callback_data="stars_grp")],
    ]
    await update.message.reply_text(
        tx(lang,"updateplan_txt",
           us=USDT_STANDARD,ss=STARS_STANDARD,cs=COINS_STANDARD,
           up=USDT_PREMIUM,sp=STARS_PREMIUM,cp=COINS_PREMIUM,
           sg=STARS_GROUP,pe=pe,plan=pn,exp=exp,coins=u["coins"]),
        reply_markup=InlineKeyboardMarkup(kb), parse_mode="Markdown")

async def coins_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    kb = []
    if u["coins"] >= COINS_STANDARD:
        kb.append([InlineKeyboardButton(tx(lang,"btn_sc"),callback_data="coins_std")])
    if u["coins"] >= COINS_PREMIUM:
        kb.append([InlineKeyboardButton(tx(lang,"btn_pc"),callback_data="coins_prm")])
    await update.message.reply_text(
        tx(lang,"coins_info",coins=u["coins"],cpm=COINS_PER_MSG,
           aff=AFFILIATE_PCT,std=COINS_STANDARD,prm=COINS_PREMIUM),
        reply_markup=InlineKeyboardMarkup(kb) if kb else None,
        parse_mode="Markdown")

async def referral_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    bot = await ctx.bot.get_me()
    link = f"https://t.me/{bot.username}?start={u['referral_code']}"
    count = u["referral_count"]
    ss = "✅" if u["claimed_standard"] else ("🔓 Ready!" if count >= 10 else f"{count}/10")
    ps = "✅" if u["claimed_premium"] else ("🔓 Ready!" if count >= 30 else f"{count}/30")
    kb = []
    if count >= 10 and not u["claimed_standard"]:
        kb.append([InlineKeyboardButton(tx(lang,"claim_std"),callback_data="claim_std")])
    if count >= 30 and not u["claimed_premium"]:
        kb.append([InlineKeyboardButton(tx(lang,"claim_prm"),callback_data="claim_prm")])
    await update.message.reply_text(
        tx(lang,"referral_info",link=link,count=count,ss=ss,ps=ps,aff=AFFILIATE_PCT),
        reply_markup=InlineKeyboardMarkup(kb) if kb else None,
        parse_mode="Markdown")

async def affiliate_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    bot = await ctx.bot.get_me()
    link = f"https://t.me/{bot.username}?start={u['affiliate_code'] or ''}"
    await update.message.reply_text(
        tx(lang,"aff_info",link=link,count=u["referral_count"],
           earned=u["affiliate_earnings"],aff=AFFILIATE_PCT),
        parse_mode="Markdown")

async def gift_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    today = date.today()
    if u["last_gift_claim"] and u["last_gift_claim"] >= today:
        await update.message.reply_text(tx(lang,"gift_done"),parse_mode="Markdown"); return
    streak = u["streak_count"]
    streak = streak+1 if (u["last_gift_claim"] and (today-u["last_gift_claim"]).days==1) else 1
    try:
        conn = db(); c = conn.cursor()
        c.execute("""UPDATE users SET streak_count=%s,last_gift_claim=%s,
            usage_chat=usage_chat+5,coins=coins+25 WHERE user_id=%s""",
                  (streak, today, user.id))
        conn.commit(); conn.close()
    except: pass
    if streak >= 10 and not u["streak_claimed"]:
        set_plan(user.id, "standard", 10)
        try:
            conn = db(); c = conn.cursor()
            c.execute("UPDATE users SET streak_claimed=TRUE WHERE user_id=%s", (user.id,))
            conn.commit(); conn.close()
        except: pass
        await update.message.reply_text(tx(lang,"gift_streak"),parse_mode="Markdown")
    else:
        await update.message.reply_text(tx(lang,"gift_ok",streak=streak),parse_mode="Markdown")

async def top_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    try:
        conn = db(); c = conn.cursor()
        c.execute("""SELECT username,full_name,total_messages,plan,coins
            FROM users WHERE is_blocked=FALSE
            ORDER BY total_messages DESC LIMIT 10""")
        rows = c.fetchall(); conn.close()
        medals = ["🥇","🥈","🥉","4️⃣","5️⃣","6️⃣","7️⃣","8️⃣","9️⃣","🔟"]
        pem = {"free":"🆓","standard":"⭐","premium":"💎"}
        text = tx(lang, "top_title")
        for i, (un, fn, msgs, plan, coins) in enumerate(rows):
            name = f"@{un}" if un else (fn or "Anonymous")
            text += f"{medals[i]} {name} {pem.get(plan,'🆓')} — *{msgs:,}* | 🪙{coins:,}\n"
        await update.message.reply_text(text, parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def stats_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    try:
        conn = db(); c = conn.cursor()
        c.execute("""SELECT usage_chat,usage_search,usage_image,usage_post,
            usage_biznes,usage_pdf,usage_cv,usage_email,usage_tts
            FROM users WHERE user_id=%s""", (user.id,))
        row = c.fetchone(); conn.close()
        if not row: await update.message.reply_text("No stats yet."); return
        chat,search,image,post,biznes,pdf,cv,email,tts = row
        pe, pn = plan_info(u); exp = exp_str(u)
        await update.message.reply_text(
            tx(lang,"stats_info",pe=pe,plan=pn,exp=exp,coins=u["coins"],
               refs=u["referral_count"],streak=u["streak_count"],
               total=u["total_messages"],chat=chat,search=search,
               image=image,pdf=pdf,cv=cv,email=email,tts=tts,post=post),
            parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def promo_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if not ctx.args:
        await update.message.reply_text(tx(lang,"promo_enter"),parse_mode="Markdown"); return
    code = ctx.args[0].upper()
    try:
        conn = db(); c = conn.cursor()
        c.execute("""SELECT reward_type,reward_value,max_uses,used_count,expires_at
            FROM promo_codes WHERE code=%s""", (code,))
        promo = c.fetchone()
        if not promo:
            await update.message.reply_text(tx(lang,"promo_bad"),parse_mode="Markdown")
            conn.close(); return
        rtype, rval, maxu, used, exp = promo
        if (exp and datetime.now() > exp) or used >= maxu:
            await update.message.reply_text(tx(lang,"promo_bad"),parse_mode="Markdown")
            conn.close(); return
        c.execute("SELECT 1 FROM promo_uses WHERE user_id=%s AND code=%s", (user.id, code))
        if c.fetchone():
            await update.message.reply_text(tx(lang,"promo_used"),parse_mode="Markdown")
            conn.close(); return
        c.execute("INSERT INTO promo_uses(user_id,code) VALUES(%s,%s)", (user.id, code))
        c.execute("UPDATE promo_codes SET used_count=used_count+1 WHERE code=%s", (code,))
        conn.commit(); conn.close()
        if rtype == "standard": set_plan(user.id,"standard",rval); reward=f"⭐ Standard {rval} days!"
        elif rtype == "premium": set_plan(user.id,"premium",rval); reward=f"💎 Premium {rval} days!"
        elif rtype == "coins": add_coins(user.id,rval); reward=f"🪙 +{rval:,} coins!"
        else: reward = "✅"
        await update.message.reply_text(tx(lang,"promo_ok",reward=reward),parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def createpromo_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Access denied!"); return
    if len(ctx.args) < 3:
        await update.message.reply_text(
            "📋 *Create Promo Code*\n\n"
            "`/createpromo [type] [value] [max_uses] [custom_code]`\n\n"
            "Types: `standard` `premium` `coins`\n\n"
            "Examples:\n"
            "`/createpromo standard 7 50` — auto code\n"
            "`/createpromo premium 30 10 VIP2025` — custom code\n"
            "`/createpromo coins 5000 100 BONUS50`",
            parse_mode="Markdown"); return
    try:
        rtype, rval, maxu = ctx.args[0], int(ctx.args[1]), int(ctx.args[2])
        code = ctx.args[3].upper() if len(ctx.args) > 3 else rnd(10)
        conn = db(); c = conn.cursor()
        c.execute("""INSERT INTO promo_codes(code,reward_type,reward_value,max_uses,expires_at)
            VALUES(%s,%s,%s,%s,%s)""",
            (code, rtype, rval, maxu, datetime.now()+timedelta(days=30)))
        conn.commit(); conn.close()
        await update.message.reply_text(
            f"✅ *Promo Code Created!*\n━━━━━━━━━━━━━━━━━━━━\n\n"
            f"🎟️ Code: `{code}`\n"
            f"📦 Type: *{rtype}*\n"
            f"💰 Value: *{rval}*\n"
            f"👥 Max uses: *{maxu}*\n"
            f"⏰ Expires: 30 days\n\n"
            f"Share this code with your users!",
            parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def group_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    # Admin uchun bepul yoqish tugmasi
    if u["is_admin"]:
        kb = [[InlineKeyboardButton("✅ Activate for Free (Admin)", callback_data="admin_activate_group")]]
    else:
        kb = [[InlineKeyboardButton(tx(lang,"btn_grp"),callback_data="stars_grp")]]
    await update.message.reply_text(
        tx(lang,"group_info",s=STARS_GROUP),
        reply_markup=InlineKeyboardMarkup(kb),
        parse_mode="Markdown")

async def activate_group_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    chat = update.effective_chat
    if chat.type == "private":
        await update.message.reply_text("⚠️ This command works only in groups!"); return
    user = update.effective_user
    u = get_user(user.id); lang = u["language"]
    activate_group(chat.id, user.id)
    await update.message.reply_text(tx(lang,"group_ok"),parse_mode="Markdown")

async def imagine_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    prompt = " ".join(ctx.args)
    if not prompt:
        await update.message.reply_text(tx(lang,"ex_imagine"),parse_mode="Markdown"); return
    if limits["imagine"] == 0:
        kb = [
            [InlineKeyboardButton(tx(lang,"btn_uprm"),callback_data="card_prm")],
            [InlineKeyboardButton(tx(lang,"btn_img"),callback_data=f"img_pay_{prompt[:50]}")]
        ]
        await update.message.reply_text(
            tx(lang,"locked_img",s=STARS_IMAGINE),
            reply_markup=InlineKeyboardMarkup(kb),
            parse_mode="Markdown"); return
    if not check_limit(user.id,"imagine",limits["imagine"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"gen_image"),parse_mode="Markdown")
    img = gen_image(prompt)
    if img:
        buf = io.BytesIO(img); buf.name = "image.png"
        await update.message.reply_photo(photo=buf, caption=f"🎨 {prompt}")
    else:
        await update.message.reply_text("❌ Could not generate image. Try a different prompt.")

async def translate_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    limits = get_limits(u["plan"], u["is_admin"])
    if not check_limit(user.id,"translate",limits["translate"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    text = " ".join(ctx.args)
    if not text:
        await update.message.reply_text(tx(lang,"ex_translate"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"translating"),parse_mode="Markdown")
    try:
        reply = await ai_once(f"Translate the following text accurately, keeping the meaning and tone:\n{text}")
        await update.message.reply_text(f"🌐 {reply}")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def code_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    limits = get_limits(u["plan"], u["is_admin"])
    if not check_limit(user.id,"code",limits["code"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    text = " ".join(ctx.args)
    if not text:
        await update.message.reply_text(tx(lang,"ex_code"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"writing_code"),parse_mode="Markdown")
    try:
        reply = await ai_once(
            f"Write clean, well-commented, production-ready code for: {text}\n"
            f"Include explanation of how it works.")
        await update.message.reply_text(f"💻 {reply}")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def document_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if limits["document"] == 0:
        kb = [[InlineKeyboardButton(tx(lang,"btn_ustd"),callback_data="card_std")]]
        await update.message.reply_text(
            tx(lang,"locked_std"),reply_markup=InlineKeyboardMarkup(kb),parse_mode="Markdown"); return
    if not check_limit(user.id,"document",limits["document"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    text = " ".join(ctx.args)
    if not text:
        await update.message.reply_text(tx(lang,"ex_document"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"writing_doc"),parse_mode="Markdown")
    try:
        reply = await ai_once(
            f"Create a professional, complete document for: {text}\n"
            f"Format it properly with all necessary sections. Use same language as the request.")
        await update.message.reply_text(f"📋 {reply}")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def ai_sound_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if limits["tts"] == 0:
        kb = [[InlineKeyboardButton(tx(lang,"btn_ustd"),callback_data="card_std")]]
        await update.message.reply_text(
            tx(lang,"locked_std"),reply_markup=InlineKeyboardMarkup(kb),parse_mode="Markdown"); return
    if not check_limit(user.id,"tts",limits["tts"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    text = " ".join(ctx.args)
    if not text:
        await update.message.reply_text(tx(lang,"ex_voice"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"gen_voice"),parse_mode="Markdown")
    try:
        tts_lang = detect_gtts(text)
        tts = gTTS(text=text, lang=tts_lang)
        buf = io.BytesIO(); tts.write_to_fp(buf); buf.seek(0)
        await update.message.reply_voice(voice=buf)
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def pptx_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if limits["pptx"] == 0:
        kb = [[InlineKeyboardButton(tx(lang,"btn_uprm"),callback_data="card_prm")]]
        await update.message.reply_text(
            tx(lang,"locked_prm"),reply_markup=InlineKeyboardMarkup(kb),parse_mode="Markdown"); return
    topic = " ".join(ctx.args)
    if not topic:
        await update.message.reply_text(tx(lang,"ex_pptx"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"creating_pptx",tp=topic),parse_mode="Markdown")
    try:
        raw = await ai_once(
            f"Create a presentation about: '{topic}'.\n"
            f"Reply with ONLY valid JSON, nothing else:\n"
            f'{{"title":"presentation title","slides":[{{"title":"slide title","points":["point 1","point 2","point 3"]}}]}}\n'
            f"Create minimum 6 slides. Use the same language as the topic.")
        raw = raw.strip()
        if "```" in raw:
            raw = raw.split("```")[1]
            if raw.startswith("json"): raw = raw[4:]
        data = json.loads(raw)
        path = make_pptx(data["title"], data["slides"])
        with open(path, "rb") as f:
            await update.message.reply_document(
                f, filename=f"{topic}.pptx", caption=tx(lang,"ready"))
        os.remove(path)
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def word_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if limits["word"] == 0:
        kb = [[InlineKeyboardButton(tx(lang,"btn_uprm"),callback_data="card_prm")]]
        await update.message.reply_text(
            tx(lang,"locked_prm"),reply_markup=InlineKeyboardMarkup(kb),parse_mode="Markdown"); return
    topic = " ".join(ctx.args)
    if not topic:
        await update.message.reply_text(tx(lang,"ex_word"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"creating_word",tp=topic),parse_mode="Markdown")
    try:
        raw = await ai_once(
            f"Create a Word document about: '{topic}'.\n"
            f"Reply with ONLY valid JSON, nothing else:\n"
            f'{{"title":"document title","sections":[{{"title":"section title","points":["point 1","point 2","point 3"]}}]}}\n'
            f"Create minimum 5 sections. Use same language as topic.")
        raw = raw.strip()
        if "```" in raw:
            raw = raw.split("```")[1]
            if raw.startswith("json"): raw = raw[4:]
        data = json.loads(raw)
        path = make_docx(data["title"], data["sections"])
        with open(path, "rb") as f:
            await update.message.reply_document(
                f, filename=f"{topic}.docx", caption=tx(lang,"ready"))
        os.remove(path)
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def cv_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if limits["cv"] == 0:
        kb = [[InlineKeyboardButton(tx(lang,"btn_ustd"),callback_data="card_std")]]
        await update.message.reply_text(
            tx(lang,"locked_std"),reply_markup=InlineKeyboardMarkup(kb),parse_mode="Markdown"); return
    if not check_limit(user.id,"cv",limits["cv"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    info = " ".join(ctx.args)
    if not info:
        await update.message.reply_text(tx(lang,"ex_cv"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"writing_cv"),parse_mode="Markdown")
    try:
        reply = await ai_once(
            f"Write a professional, ATS-friendly CV/Resume for: {info}\n"
            f"Include: Professional Summary, Work Experience, Technical Skills, "
            f"Education, Key Achievements. Make it impressive and detailed.")
        await update.message.reply_text(f"👤 *CV:*\n\n{reply}", parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def email_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if limits["email"] == 0:
        kb = [[InlineKeyboardButton(tx(lang,"btn_ustd"),callback_data="card_std")]]
        await update.message.reply_text(
            tx(lang,"locked_std"),reply_markup=InlineKeyboardMarkup(kb),parse_mode="Markdown"); return
    if not check_limit(user.id,"email",limits["email"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    topic = " ".join(ctx.args)
    if not topic:
        await update.message.reply_text(tx(lang,"ex_email"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"writing_email"),parse_mode="Markdown")
    try:
        reply = await ai_once(
            f"Write a professional email about: {topic}\n"
            f"Include: Subject line, proper greeting, clear body, professional closing. "
            f"Use same language as the request.")
        await update.message.reply_text(f"📧 *Email:*\n\n{reply}", parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def post_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if not check_limit(user.id,"post",limits["post"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    topic = " ".join(ctx.args)
    if not topic:
        await update.message.reply_text(tx(lang,"ex_post"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"writing_post"),parse_mode="Markdown")
    try:
        reply = await ai_once(
            f"Write an engaging, viral social media post about: {topic}\n"
            f"Include: Strong hook, valuable content, clear call-to-action, relevant hashtags. "
            f"Use same language as the request.")
        await update.message.reply_text(f"📱 *Post:*\n\n{reply}", parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def biznes_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if not check_limit(user.id,"biznes",limits["biznes"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    idea = " ".join(ctx.args)
    if not idea:
        await update.message.reply_text(tx(lang,"ex_biznes"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"writing_biz"),parse_mode="Markdown")
    try:
        reply = await ai_once(
            f"Write a comprehensive business plan for: {idea}\n"
            f"Sections: Executive Summary, Market Analysis, Products/Services, "
            f"Marketing Strategy, Financial Projections, Competitive Advantage. "
            f"Be detailed and professional.")
        await update.message.reply_text(f"💼 *Business Plan:*\n\n{reply}", parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def admin_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Access denied!"); return
    try:
        conn = db(); c = conn.cursor()
        c.execute("SELECT COUNT(*) FROM users"); total = c.fetchone()[0]
        c.execute("SELECT COUNT(*) FROM users WHERE plan='standard'"); std = c.fetchone()[0]
        c.execute("SELECT COUNT(*) FROM users WHERE plan='premium'"); prm = c.fetchone()[0]
        c.execute("SELECT COUNT(*) FROM users WHERE is_blocked=TRUE"); blk = c.fetchone()[0]
        c.execute("SELECT SUM(coins) FROM users"); tc = c.fetchone()[0] or 0
        c.execute("SELECT COUNT(*) FROM group_modes"); grps = c.fetchone()[0]
        conn.close()
        rev = std * USDT_STANDARD + prm * USDT_PREMIUM
    except: total=std=prm=blk=tc=grps=rev=0
    await update.message.reply_text(
        f"🔧 *Admin Panel*\n━━━━━━━━━━━━━━━━━━━━\n\n"
        f"👥 Total users: *{total}*\n"
        f"🆓 Free: *{total-std-prm}*\n"
        f"⭐ Standard: *{std}* | 💎 Premium: *{prm}*\n"
        f"🚫 Blocked: *{blk}* | 👥 Groups: *{grps}*\n"
        f"💰 Est. Revenue: ~*${rev}*\n"
        f"🪙 Total coins: *{tc:,}*\n\n"
        f"📋 *Commands:*\n"
        f"`/users` — All users list\n"
        f"`/find [@username or id]` — Find & manage user\n"
        f"`/broadcast [message]` — Message all users\n"
        f"`/createpromo [type] [val] [uses] [code?]` — Create promo\n"
        f"`/addcoins [user_id] [amount]` — Add coins\n"
        f"`/maintenance` — Send maintenance alert",
        parse_mode="Markdown")

async def maintenance_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Access denied!"); return
    await update.message.reply_text("⏳ Sending maintenance alert to all users...")
    await broadcast_maintenance(ctx.application)
    await update.message.reply_text("✅ Maintenance alert sent to all users!")

async def users_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Access denied!"); return
    try:
        conn = db(); c = conn.cursor()
        c.execute("""SELECT user_id,username,plan,is_blocked FROM users
            ORDER BY CASE plan WHEN 'premium' THEN 1 WHEN 'standard' THEN 2 ELSE 3 END,
            joined_at DESC""")
        rows = c.fetchall(); conn.close()
        if not rows: await update.message.reply_text("No users."); return
        pem = {"free":"🆓","standard":"⭐","premium":"💎"}
        text = f"👥 *All Users: {len(rows)}*\n━━━━━━━━━━━━━━━━━━━━\n\n"
        for uid, un, plan, blk in rows:
            n = f"@{un}" if un else "no_username"
            b = " 🚫" if blk else ""
            text += f"{pem.get(plan,'🆓')} {n} `{uid}`{b}\n"
        if len(text) > 4000:
            for p in [text[i:i+4000] for i in range(0, len(text), 4000)]:
                await update.message.reply_text(p, parse_mode="Markdown")
        else:
            await update.message.reply_text(text, parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def find_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Access denied!"); return
    if not ctx.args:
        await update.message.reply_text("Usage: /find [id or @username]"); return
    try:
        arg = ctx.args[0]; conn = db(); c = conn.cursor()
        if arg.startswith("@"):
            c.execute("""SELECT plan,expires_at,is_blocked,full_name,username,
                user_id,coins,total_messages,warning_count
                FROM users WHERE username=%s""", (arg[1:],))
        else:
            c.execute("""SELECT plan,expires_at,is_blocked,full_name,username,
                user_id,coins,total_messages,warning_count
                FROM users WHERE user_id=%s""", (int(arg),))
        row = c.fetchone(); conn.close()
        if not row: await update.message.reply_text(f"❌ Not found: {arg}"); return
        plan, exp, blk, fn, un, tid, coins, msgs, warns = row
        pem = {"free":"🆓","standard":"⭐","premium":"💎"}
        exp_s = exp.strftime("%d.%m.%Y") if exp else "—"
        kb = [
            [InlineKeyboardButton("🆓 Free",callback_data=f"ap_free_{tid}"),
             InlineKeyboardButton("⭐ Std",callback_data=f"ap_std_{tid}"),
             InlineKeyboardButton("💎 Prm",callback_data=f"ap_prm_{tid}")],
            [InlineKeyboardButton(
                "✅ Unblock" if blk else "🚫 Block",
                callback_data=f"ap_unblock_{tid}" if blk else f"ap_block_{tid}")]
        ]
        await update.message.reply_text(
            f"👤 *User Info*\n━━━━━━━━━━━━━━━━━━━━\n"
            f"🆔 ID: `{tid}`\n"
            f"👤 Name: {fn or '—'} | @{un or '—'}\n"
            f"{pem.get(plan,'🆓')} Plan: *{plan.upper()}* | 📅 {exp_s}\n"
            f"🪙 Coins: {coins:,} | 💬 Msgs: {msgs:,} | ⚠️ Warns: {warns}\n"
            f"{'🚫 BLOCKED' if blk else '✅ Active'}",
            reply_markup=InlineKeyboardMarkup(kb), parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def broadcast_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Access denied!"); return
    msg = " ".join(ctx.args)
    if not msg: await update.message.reply_text("Usage: /broadcast [message]"); return
    try:
        conn = db(); c = conn.cursor()
        c.execute("SELECT user_id FROM users WHERE is_blocked=FALSE")
        rows = c.fetchall(); conn.close()
        sent = 0
        for row in rows:
            try: await ctx.bot.send_message(chat_id=row[0], text=f"📢 {msg}"); sent += 1
            except: pass
        await update.message.reply_text(f"✅ Sent to *{sent}* users!", parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def addcoins_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Access denied!"); return
    if len(ctx.args) < 2:
        await update.message.reply_text("Usage: /addcoins [user_id] [amount]"); return
    try:
        uid, amount = int(ctx.args[0]), int(ctx.args[1])
        add_coins(uid, amount)
        await update.message.reply_text(
            f"✅ Added *{amount:,}* coins to `{uid}`", parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def precheckout(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    await update.pre_checkout_query.answer(ok=True)

async def payment_success(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    payload = update.message.successful_payment.invoice_payload
    u = get_user(user.id); lang = u["language"]
    if payload == "std_stars":
        set_plan(user.id, "standard", 30)
        await update.message.reply_text("✅ ⭐ *Standard* activated for 30 days! 🎉", parse_mode="Markdown")
    elif payload == "prm_stars":
        set_plan(user.id, "premium", 30)
        await update.message.reply_text("✅ 💎 *Premium* activated for 30 days! 🎉", parse_mode="Markdown")
    elif payload == "grp_stars":
        await update.message.reply_text(
            "✅ *Group mode purchased!*\n\nNow add me to your group and send /activate\\_group",
            parse_mode="Markdown")
    elif payload.startswith("img_"):
        prompt = payload[4:]
        await update.message.reply_text(tx(lang,"gen_image"),parse_mode="Markdown")
        img = gen_image(prompt)
        if img:
            buf = io.BytesIO(img); buf.name = "image.png"
            await update.message.reply_photo(photo=buf, caption=f"🎨 {prompt}")
        else:
            await update.message.reply_text("❌ Could not generate. Try again.")

async def callback(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query; await q.answer()
    data = q.data; uid = q.from_user.id
    u = get_user(uid); lang = u["language"]

    if data.startswith("lang_"):
        nl = data[5:]
        lca = u["language_changed_at"]
        if lca and datetime.now() - lca < timedelta(hours=24):
            await q.answer(tx(lang,"lang_cooldown"), show_alert=True); return
        set_language(uid, nl)
        flag, name = LANGUAGES.get(nl, ("🌐","Unknown"))
        await q.message.edit_text(f"✅ Language changed to *{flag} {name}!*", parse_mode="Markdown")

    elif data == "admin_activate_group":
        # Admin uchun bepul guruh yoqish
        if uid != ADMIN_ID:
            await q.answer("❌ Access denied!", show_alert=True); return
        await q.message.reply_text(
            "✅ *Admin Group Activation*\n\n"
            "Add me to your group, then send /activate\\_group in the group.",
            parse_mode="Markdown")

    elif data == "stars_std":
        await ctx.bot.send_invoice(
            chat_id=uid, title="⭐ Standard Plan",
            description="Standard plan for 30 days", payload="std_stars",
            currency="XTR", prices=[LabeledPrice("Standard 30 days", STARS_STANDARD)])

    elif data == "stars_prm":
        await ctx.bot.send_invoice(
            chat_id=uid, title="💎 Premium Plan",
            description="Premium plan for 30 days", payload="prm_stars",
            currency="XTR", prices=[LabeledPrice("Premium 30 days", STARS_PREMIUM)])

    elif data == "stars_grp":
        # Admin uchun bepul — Stars to'lamasdan yoqiladi
        if uid == ADMIN_ID:
            await q.message.reply_text(
                "👑 *Admin Group Mode*\n\n"
                "As admin, group mode is FREE for you!\n"
                "Add bot to your group → send /activate\\_group",
                parse_mode="Markdown")
        else:
            await ctx.bot.send_invoice(
                chat_id=uid, title="👥 Group Mode",
                description="Activate bot for your entire group", payload="grp_stars",
                currency="XTR", prices=[LabeledPrice("Group Mode", STARS_GROUP)])

    elif data.startswith("img_pay_"):
        prompt = data[8:]
        await ctx.bot.send_invoice(
            chat_id=uid, title="🎨 AI Image",
            description=f"Generate: {prompt}", payload=f"img_{prompt}",
            currency="XTR", prices=[LabeledPrice("AI Image", STARS_IMAGINE)])

    elif data == "coins_std":
        if u["coins"] >= COINS_STANDARD:
            conn = db(); c = conn.cursor()
            c.execute("UPDATE users SET coins=coins-%s WHERE user_id=%s", (COINS_STANDARD, uid))
            conn.commit(); conn.close()
            set_plan(uid, "standard", 30)
            await q.message.edit_text(tx(lang,"coins_ok_std"), parse_mode="Markdown")
        else:
            await q.answer(tx(lang,"not_enough",need=COINS_STANDARD-u["coins"]), show_alert=True)

    elif data == "coins_prm":
        if u["coins"] >= COINS_PREMIUM:
            conn = db(); c = conn.cursor()
            c.execute("UPDATE users SET coins=coins-%s WHERE user_id=%s", (COINS_PREMIUM, uid))
            conn.commit(); conn.close()
            set_plan(uid, "premium", 30)
            await q.message.edit_text(tx(lang,"coins_ok_prm"), parse_mode="Markdown")
        else:
            await q.answer(tx(lang,"not_enough",need=COINS_PREMIUM-u["coins"]), show_alert=True)

    elif data == "card_std":
        # Admin uchun bepul
        if uid == ADMIN_ID:
            set_plan(uid, "standard", 36500)  # 100 yil
            await q.message.reply_text("👑 Admin: Standard activated for free!", parse_mode="Markdown")
        else:
            kb = [[InlineKeyboardButton(tx(lang,"btn_contact"), url=f"https://t.me/{ADMIN_USERNAME}")]]
            await q.message.reply_text(
                tx(lang,"pay_card",card=CARD_NUMBER,amt=USDT_STANDARD,admin=ADMIN_USERNAME),
                reply_markup=InlineKeyboardMarkup(kb), parse_mode="Markdown")

    elif data == "card_prm":
        # Admin uchun bepul
        if uid == ADMIN_ID:
            set_plan(uid, "premium", 36500)  # 100 yil
            await q.message.reply_text("👑 Admin: Premium activated for free!", parse_mode="Markdown")
        else:
            kb = [[InlineKeyboardButton(tx(lang,"btn_contact"), url=f"https://t.me/{ADMIN_USERNAME}")]]
            await q.message.reply_text(
                tx(lang,"pay_card",card=CARD_NUMBER,amt=USDT_PREMIUM,admin=ADMIN_USERNAME),
                reply_markup=InlineKeyboardMarkup(kb), parse_mode="Markdown")

    elif data == "claim_std":
        if u["referral_count"] >= 10 and not u["claimed_standard"]:
            set_plan(uid, "standard", 15)
            conn = db(); c = conn.cursor()
            c.execute("UPDATE users SET claimed_standard=TRUE WHERE user_id=%s", (uid,))
            conn.commit(); conn.close()
            await q.message.edit_text(tx(lang,"claimed_std"), parse_mode="Markdown")
        else:
            await q.answer(tx(lang,"claim_err"), show_alert=True)

    elif data == "claim_prm":
        if u["referral_count"] >= 30 and not u["claimed_premium"]:
            set_plan(uid, "premium", 15)
            conn = db(); c = conn.cursor()
            c.execute("UPDATE users SET claimed_premium=TRUE WHERE user_id=%s", (uid,))
            conn.commit(); conn.close()
            await q.message.edit_text(tx(lang,"claimed_prm"), parse_mode="Markdown")
        else:
            await q.answer(tx(lang,"claim_err"), show_alert=True)

    elif data.startswith("ap_free_") or data.startswith("ap_std_") or data.startswith("ap_prm_"):
        parts = data.split("_")
        pm = {"free":"free","std":"standard","prm":"premium"}
        plan_key = parts[1]; tid = int(parts[2])
        plan = pm[plan_key]; days = 30 if plan != "free" else None
        set_plan(tid, plan, days)
        pem = {"free":"🆓","standard":"⭐","premium":"💎"}
        await q.message.edit_text(
            f"✅ User `{tid}` → {pem[plan]} *{plan.upper()}*", parse_mode="Markdown")
        try:
            await ctx.bot.send_message(
                tid, f"{pem[plan]} Your plan updated to *{plan.upper()}*!",
                parse_mode="Markdown")
        except: pass

    elif data.startswith("ap_block_"):
        tid = int(data.split("_")[2]); set_blocked(tid, True)
        await q.message.edit_text(f"🚫 User `{tid}` blocked!", parse_mode="Markdown")

    elif data.startswith("ap_unblock_"):
        tid = int(data.split("_")[2]); set_blocked(tid, False)
        await q.message.edit_text(f"✅ User `{tid}` unblocked!", parse_mode="Markdown")

async def handle_text(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    chat = update.effective_chat
    text = update.message.text or ""

    # Guruh chati
    if chat.type in ["group", "supergroup"]:
        if not text.lower().startswith("bot "): return
        if not is_group_active(chat.id): return
        question = text[4:].strip()
        if not question: return
        try:
            await ctx.bot.send_chat_action(chat_id=chat.id, action="typing")
            reply = await ai_chat(user.id, question)
            await update.message.reply_text(reply)
        except Exception as e:
            await update.message.reply_text(f"Error: {e}")
        return

    # Shaxsiy chat
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]

    if u["is_blocked"]:
        await update.message.reply_text(
            tx(lang,"blocked",admin=ADMIN_USERNAME), parse_mode="Markdown"); return

    if is_bad(text):
        wc = add_warning(user.id)
        if wc >= MAX_WARNINGS:
            set_blocked(user.id, True)
            await update.message.reply_text(tx(lang,"banned"),parse_mode="Markdown")
        else:
            await update.message.reply_text(
                tx(lang,"warning",c=wc,m=MAX_WARNINGS), parse_mode="Markdown")
        return

    limits = get_limits(u["plan"], u["is_admin"])
    if not u["is_admin"] and not check_limit(user.id,"chat",limits["chat"]):
        kb = [[InlineKeyboardButton(tx(lang,"btn_ustd"),callback_data="card_std")]]
        await update.message.reply_text(
            tx(lang,"limit_reached"),
            reply_markup=InlineKeyboardMarkup(kb),
            parse_mode="Markdown"); return

    total = inc_msg(user.id, u["is_admin"])
    if total == 100000 and not u["achievement_100k"] and not u["is_admin"]:
        set_plan(user.id, "standard", 5)
        try:
            conn = db(); c = conn.cursor()
            c.execute("UPDATE users SET achievement_100k=TRUE WHERE user_id=%s", (user.id,))
            conn.commit(); conn.close()
        except: pass
        await update.message.reply_text(tx(lang,"achievement"),parse_mode="Markdown")

    await ctx.bot.send_chat_action(chat_id=chat.id, action="typing")
    try:
        reply = await ai_chat(user.id, text, do_search=True, limits=limits)
        await update.message.reply_text(reply)
        await update_mem(user.id, text, reply)
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def handle_voice(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if limits["voice"] == 0:
        kb = [[InlineKeyboardButton(tx(lang,"btn_uprm"),callback_data="card_prm")]]
        await update.message.reply_text(
            tx(lang,"locked_prm"),reply_markup=InlineKeyboardMarkup(kb),parse_mode="Markdown"); return
    await ctx.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    try:
        f = await ctx.bot.get_file(update.message.voice.file_id)
        data = requests.get(f.file_path).content
        with open("v.ogg","wb") as fp: fp.write(data)
        with open("v.ogg","rb") as fp:
            tr = ai.audio.transcriptions.create(
                file=("v.ogg", fp.read()),
                model="whisper-large-v3")
        text = tr.text
        await update.message.reply_text(
            tx(lang,"you_said") + text, parse_mode="Markdown")
        reply = await ai_chat(user.id, text)
        await update.message.reply_text(reply)
        try: os.remove("v.ogg")
        except: pass
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def handle_photo(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if not u["is_admin"] and not check_limit(user.id,"image",limits["image"]):
        kb = [[InlineKeyboardButton(tx(lang,"btn_ustd"),callback_data="card_std")]]
        await update.message.reply_text(
            tx(lang,"limit_reached"),reply_markup=InlineKeyboardMarkup(kb),parse_mode="Markdown"); return
    await ctx.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    try:
        photo = update.message.photo[-1]
        f = await ctx.bot.get_file(photo.file_id)
        img_data = requests.get(f.file_path).content
        b64 = base64.b64encode(img_data).decode()
        caption = update.message.caption or "Describe this image in detail. What do you see?"
        r = ai.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{"role":"user","content":[
                {"type":"image_url","image_url":{"url":f"data:image/jpeg;base64,{b64}"}},
                {"type":"text","text":caption}
            ]}])
        await update.message.reply_text(r.choices[0].message.content)
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def handle_document(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    ensure_user(user.id, user.username, user.full_name)
    u = get_user(user.id); lang = u["language"]
    if u["is_blocked"]: return
    limits = get_limits(u["plan"], u["is_admin"])
    if limits["pdf"] == 0:
        kb = [[InlineKeyboardButton(tx(lang,"btn_ustd"),callback_data="card_std")]]
        await update.message.reply_text(
            tx(lang,"locked_std"),reply_markup=InlineKeyboardMarkup(kb),parse_mode="Markdown"); return
    if not u["is_admin"] and not check_limit(user.id,"pdf",limits["pdf"]):
        await update.message.reply_text(tx(lang,"limit_reached"),parse_mode="Markdown"); return
    doc = update.message.document
    if not doc.file_name.lower().endswith(".pdf"):
        await update.message.reply_text(tx(lang,"pdf_only"),parse_mode="Markdown"); return
    await update.message.reply_text(tx(lang,"reading_pdf"),parse_mode="Markdown")
    await ctx.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    try:
        f = await ctx.bot.get_file(doc.file_id)
        data = requests.get(f.file_path).content
        with open("tmp.pdf","wb") as fp: fp.write(data)
        pdf = fitz.open("tmp.pdf")
        text = "".join(p.get_text() for p in pdf); pdf.close()
        try: os.remove("tmp.pdf")
        except: pass
        if len(text) > 12000: text = text[:12000] + "..."
        caption = update.message.caption or "Summarize this document comprehensively and explain all key points."
        reply = await ai_chat(user.id, f"PDF Content:\n\n{text}\n\nRequest: {caption}", do_search=False)
        await update.message.reply_text(reply)
    except Exception as e:
        await update.message.reply_text(f"Error: {e}")

async def post_init(app):
    init_db()
    await app.bot.set_my_commands([
        BotCommand("start","🚀 Start"),
        BotCommand("updateplan","💰 Plans & pricing"),
        BotCommand("weather","🌦 Weather forecast"),
        BotCommand("crypto","💰 Crypto prices"),
        BotCommand("news","📰 Latest news"),
        BotCommand("imagine","🎨 AI Image (Premium)"),
        BotCommand("pptx","📊 PowerPoint (Premium)"),
        BotCommand("word","📝 Word document (Premium)"),
        BotCommand("cv","👤 Write CV (Standard+)"),
        BotCommand("email","📧 Write email (Standard+)"),
        BotCommand("document","📋 Create document (Standard+)"),
        BotCommand("ai_sound","🔊 AI Voice (Standard+)"),
        BotCommand("translate","🌐 Translate text"),
        BotCommand("code","💻 Write code"),
        BotCommand("post","📱 Marketing post"),
        BotCommand("biznes","💼 Business plan"),
        BotCommand("coins","🪙 My coins"),
        BotCommand("referral","👥 Invite & earn"),
        BotCommand("affiliate","🤝 Affiliate program"),
        BotCommand("gift","🎁 Daily bonus"),
        BotCommand("top","🏆 Top users"),
        BotCommand("stats","📊 My statistics"),
        BotCommand("promo","🎟️ Promo code"),
        BotCommand("language","🌐 Change language"),
        BotCommand("reset","🗑️ Clear history"),
        BotCommand("help","❓ Help"),
    ])

def main():
    app = ApplicationBuilder().token(TELEGRAM_TOKEN).post_init(post_init).build()

    # User commands
    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("help", help_cmd))
    app.add_handler(CommandHandler("reset", reset_cmd))
    app.add_handler(CommandHandler("updateplan", updateplan_cmd))
    app.add_handler(CommandHandler("language", language_cmd))
    app.add_handler(CommandHandler("weather", weather_cmd))
    app.add_handler(CommandHandler("crypto", crypto_cmd))
    app.add_handler(CommandHandler("news", news_cmd))
    app.add_handler(CommandHandler("imagine", imagine_cmd))
    app.add_handler(CommandHandler("translate", translate_cmd))
    app.add_handler(CommandHandler("code", code_cmd))
    app.add_handler(CommandHandler("document", document_cmd))
    app.add_handler(CommandHandler("ai_sound", ai_sound_cmd))
    app.add_handler(CommandHandler("pptx", pptx_cmd))
    app.add_handler(CommandHandler("word", word_cmd))
    app.add_handler(CommandHandler("cv", cv_cmd))
    app.add_handler(CommandHandler("email", email_cmd))
    app.add_handler(CommandHandler("post", post_cmd))
    app.add_handler(CommandHandler("biznes", biznes_cmd))
    app.add_handler(CommandHandler("coins", coins_cmd))
    app.add_handler(CommandHandler("referral", referral_cmd))
    app.add_handler(CommandHandler("affiliate", affiliate_cmd))
    app.add_handler(CommandHandler("gift", gift_cmd))
    app.add_handler(CommandHandler("top", top_cmd))
    app.add_handler(CommandHandler("stats", stats_cmd))
    app.add_handler(CommandHandler("promo", promo_cmd))
    app.add_handler(CommandHandler("group", group_cmd))
    app.add_handler(CommandHandler("activate_group", activate_group_cmd))

    # Admin commands
    app.add_handler(CommandHandler("admin", admin_cmd))
    app.add_handler(CommandHandler("maintenance", maintenance_cmd))
    app.add_handler(CommandHandler("users", users_cmd))
    app.add_handler(CommandHandler("find", find_cmd))
    app.add_handler(CommandHandler("broadcast", broadcast_cmd))
    app.add_handler(CommandHandler("createpromo", createpromo_cmd))
    app.add_handler(CommandHandler("addcoins", addcoins_cmd))

    # Payments
    app.add_handler(PreCheckoutQueryHandler(precheckout))
    app.add_handler(MessageHandler(filters.SUCCESSFUL_PAYMENT, payment_success))

    # Callbacks & messages
    app.add_handler(CallbackQueryHandler(callback))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))
    app.add_handler(MessageHandler(filters.VOICE, handle_voice))
    app.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    app.add_handler(MessageHandler(filters.Document.PDF, handle_document))

    print(f"🚀 {BOT_NAME} is running!")
    app.run_polling()

if __name__ == "__main__":
    main()
